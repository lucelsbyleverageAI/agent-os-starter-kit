"""Generate E18 Process One-Pager PowerPoint presentations.

This tool creates professionally formatted process one-pagers from structured inputs,
with intelligent text validation, markdown support, and Supabase storage integration.
"""

import json
import os
import re
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

from ..base import CustomTool, ToolParameter
from ...config import settings
from ...utils.logging import get_logger
from ...utils.exceptions import ToolExecutionError

logger = get_logger(__name__)

# Template path
TEMPLATE_PATH = Path(__file__).parent / "assets" / "e18-process-one-pager-template.pptx"

# Word/character limits based on template analysis
LIMITS = {
    "headline_statement": {"words": 6, "tolerance": 0.2},  # ~6 words max
    "process_name": {"words": 12, "tolerance": 0.2},  # ~12 words max
    "nhs_org_name": {"chars": 50, "tolerance": 0.2},  # ~50 chars max
    "process_category": {"words": 5, "tolerance": 0.2},  # ~5 words max
    "manual_process": {"words": 70, "tolerance": 0.2},  # ~70 words max
    "challenges": {"words": 70, "tolerance": 0.2},  # ~70 words max
    "rpa_solution": {"words": 70, "tolerance": 0.2},  # ~70 words max
    "benefits_item": {"words": 8, "tolerance": 0.25},  # ~8 words per bullet point (max 10 with tolerance)
    "benefits_count": {"min": 3, "max": 5},  # 3-5 bullet points
}


class GenerateProcessOnePagerTool(CustomTool):
    """Generate E18 Process One-Pager PowerPoint presentation.

    Creates a professional, branded PowerPoint presentation following E18's
    process one-pager template. Includes intelligent text validation,
    markdown formatting support, and automated file storage.
    """

    toolkit_name = "e18_utility"
    toolkit_display_name = "E18 Utilities"

    @property
    def name(self) -> str:
        return "generate_process_one_pager"

    @property
    def description(self) -> str:
        return (
            "Generate a professional E18 Process One-Pager PowerPoint presentation. "
            "Takes structured process information (name, description, challenges, solution, benefits) "
            "and creates a formatted, branded PPTX file following E18's template design. "
            "Supports markdown formatting (bold, bullets) and automatically validates text length. "
            "Returns a secure download URL for the generated file. Ensure all text is written in British English, not US."
        )

    def get_parameters(self) -> List[ToolParameter]:
        return [
            ToolParameter(
                name="headline_statement",
                type="string",
                description=(
                    "Bold, impactful headline statement (6-7 words max). "
                    "Keep concise and action-oriented. Use outcome metrics where possible."
                ),
                required=True
            ),
            ToolParameter(
                name="process_name",
                type="string",
                description=(
                    "Descriptive name of the RPA process (3-6 words max). "
                    "Example: 'Waiting List Validation for Outpatient Services'. "
                    "Should clearly identify what the process does."
                ),
                required=True
            ),
            ToolParameter(
                name="nhs_org_name",
                type="string",
                description=(
                    "NHS organisation name. Full name is preferred to abbreviations. "
                    "Example: 'NHS Greater Manchester'. 'Leeds Teaching Hospitals NHS Trust'."
                ),
                required=True
            ),
            ToolParameter(
                name="process_category",
                type="string",
                description=(
                    "Process category or department (3-5 words max). "
                    "Examples: 'Outpatients', 'Maternity', 'Cancer', . "
                    "Indicates which clinical area this process belongs to."
                ),
                required=True
            ),
            ToolParameter(
                name="systems_used",
                type="array",
                description=(
                    "List of 1-5 systems/applications used in the process. "
                    "Examples: ['PAS', 'RADAR', 'Patient Hub'], ['Lorenzo', 'EPR', 'PACS']. "
                    "If you don't know the systems used, return 'Information not available yet' as a single string."
                    "Will be displayed with pipe separators (e.g., 'PAS | RADAR | Patient Hub')."
                ),
                required=True,
                items={"type": "string"}
            ),
            ToolParameter(
                name="manual_process",
                type="string",
                description=(
                    "Description of the manual process before automation (60-70 words max). "
                    "Explain what staff had to do manually, the steps involved, and time required. "
                    "Primarily prose, but can use **bold** for emphasis. "
                    "Example: 'Staff manually reviewed waiting lists daily, cross-referencing "
                    "patient records across **three separate systems**. This involved logging "
                    "into each system, extracting data, and manually validating appointment dates.'"
                ),
                required=True
            ),
            ToolParameter(
                name="challenges",
                type="string",
                description=(
                    "Key challenges with the manual process (60-70 words max). "
                    "Focus on pain points, errors, inefficiencies, and impacts. "
                    "Primarily prose, but can use **bold** for emphasis. "
                    "Example: 'The manual process was **time-consuming** (4 hours daily) and "
                    "**error-prone** due to data entry mistakes. Staff reported high frustration "
                    "levels, and validation delays led to appointment scheduling issues.'"
                ),
                required=True
            ),
            ToolParameter(
                name="rpa_solution",
                type="string",
                description=(
                    "Description of the RPA solution implemented (60-70 words max). "
                    "Explain what the automation does, which platform/tools, and key features. "
                    "Primarily prose, but can use **bold** for emphasis. "
                    "Example: 'Implemented **Blue Prism RPA bot** that automatically logs into "
                    "all three systems nightly, extracts patient data, performs validation checks "
                    "against business rules, and generates exception reports for staff review.'"
                ),
                required=True
            ),
            ToolParameter(
                name="benefits",
                type="array",
                description=(
                    "List of 4-5 quantifiable benefits as concise bullet points (~8 words each). "
                    "Each item should be a string with measurable outcomes. "
                    "Use **bold** markdown for key metrics within each bullet point. "
                    "Focus on: time saved, accuracy improved, cost reduction, staff impact. "
                    "Example: ['**35,000 cases processed** annually with full automation', "
                    "'**£53,000 annual cost savings** through reduced manual workload', "
                    "'**85% reduction** in manual processing time', "
                    "'**Significant improvement in data accuracy** reducing clinical risks', "
                    "'**Increased staff capacity** to focus on patient-facing activities']"
                ),
                required=True,
                items={"type": "string"}
            ),
        ]

    async def _execute_impl(self, user_id: str, **kwargs: Any) -> str:
        """Generate PowerPoint presentation from inputs."""
        try:
            # Extract parameters
            headline_statement = kwargs.get("headline_statement", "")
            process_name = kwargs.get("process_name", "")
            nhs_org_name = kwargs.get("nhs_org_name", "")
            process_category = kwargs.get("process_category", "")
            systems_used = kwargs.get("systems_used", [])
            manual_process = kwargs.get("manual_process", "")
            challenges = kwargs.get("challenges", "")
            rpa_solution = kwargs.get("rpa_solution", "")
            benefits = kwargs.get("benefits", "")

            # Validate inputs
            validation_errors = self._validate_inputs(
                headline_statement=headline_statement,
                process_name=process_name,
                nhs_org_name=nhs_org_name,
                process_category=process_category,
                manual_process=manual_process,
                challenges=challenges,
                rpa_solution=rpa_solution,
                benefits=benefits
            )

            if validation_errors:
                raise ToolExecutionError(
                    "generate_process_one_pager",
                    f"Input validation failed:\n" + "\n".join(validation_errors)
                )

            # Load template
            if not TEMPLATE_PATH.exists():
                raise ToolExecutionError(
                    "generate_process_one_pager",
                    f"Template file not found: {TEMPLATE_PATH}"
                )

            prs = Presentation(str(TEMPLATE_PATH))

            # Get the first (and only) slide
            if len(prs.slides) == 0:
                raise ToolExecutionError(
                    "generate_process_one_pager",
                    "Template has no slides"
                )

            slide = prs.slides[0]

            # Replace placeholders in template
            self._populate_template(
                slide=slide,
                headline_statement=headline_statement,
                process_name=process_name,
                nhs_org_name=nhs_org_name,
                process_category=process_category,
                systems_used=systems_used,
                manual_process=manual_process,
                challenges=challenges,
                rpa_solution=rpa_solution,
                benefits=benefits
            )

            # Save to bytes
            pptx_buffer = BytesIO()
            prs.save(pptx_buffer)
            pptx_bytes = pptx_buffer.getvalue()

            # Generate sanitized filename
            safe_process_name = re.sub(r'[^a-zA-Z0-9_-]', '_', process_name)[:50]
            timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
            filename = f"{timestamp}_{safe_process_name}.pptx"

            # Upload to Supabase storage
            download_url = await self._upload_to_storage(
                user_id=user_id,
                filename=filename,
                file_bytes=pptx_bytes
            )

            return json.dumps({
                "success": True,
                "filename": filename,
                "download_url": download_url,
                "size_bytes": len(pptx_bytes),
                "process_name": process_name,
                "message": (
                    f"Successfully generated process one-pager: {process_name}. "
                    f"Click the download URL to access the file."
                )
            }, indent=2)

        except ToolExecutionError:
            raise
        except Exception as e:
            logger.error(f"Failed to generate PowerPoint: {e}", exc_info=True)
            raise ToolExecutionError(
                "generate_process_one_pager",
                f"Unexpected error generating PowerPoint: {str(e)}"
            )

    def _validate_inputs(
        self,
        headline_statement: str,
        process_name: str,
        nhs_org_name: str,
        process_category: str,
        manual_process: str,
        challenges: str,
        rpa_solution: str,
        benefits: List[str]
    ) -> List[str]:
        """Validate input lengths against limits.

        Returns list of validation error messages. Empty list if all valid.
        """
        errors = []

        # Word-based validations for string fields
        word_fields = {
            "headline_statement": headline_statement,
            "process_name": process_name,
            "process_category": process_category,
            "manual_process": manual_process,
            "challenges": challenges,
            "rpa_solution": rpa_solution,
        }

        for field_name, text in word_fields.items():
            word_count = len(text.split())
            limit_info = LIMITS[field_name]
            max_words = limit_info["words"]
            tolerance = limit_info["tolerance"]

            # Only error if significantly over (>20%)
            if word_count > max_words * (1 + tolerance):
                errors.append(
                    f"{field_name}: {word_count} words exceeds limit of {max_words} words "
                    f"(max {int(max_words * (1 + tolerance))} with tolerance). Please shorten."
                )

        # Character-based validation for NHS org name
        char_limit = LIMITS["nhs_org_name"]["chars"]
        tolerance = LIMITS["nhs_org_name"]["tolerance"]
        if len(nhs_org_name) > char_limit * (1 + tolerance):
            errors.append(
                f"nhs_org_name: {len(nhs_org_name)} characters exceeds limit of "
                f"{char_limit} characters. Please shorten or use abbreviations."
            )

        # Validate benefits array
        count_limits = LIMITS["benefits_count"]
        if len(benefits) < count_limits["min"]:
            errors.append(
                f"benefits: Must provide at least {count_limits['min']} bullet points. "
                f"Currently have {len(benefits)}."
            )
        elif len(benefits) > count_limits["max"]:
            errors.append(
                f"benefits: Too many bullet points. Maximum is {count_limits['max']}, "
                f"currently have {len(benefits)}."
            )

        # Validate each benefit item's word count
        item_limits = LIMITS["benefits_item"]
        max_words_per_item = item_limits["words"]
        tolerance = item_limits["tolerance"]

        for i, benefit in enumerate(benefits, 1):
            word_count = len(benefit.split())
            if word_count > max_words_per_item * (1 + tolerance):
                errors.append(
                    f"benefits[{i}]: {word_count} words exceeds limit of {max_words_per_item} words "
                    f"(max {int(max_words_per_item * (1 + tolerance))} with tolerance). Please shorten."
                )

        return errors

    def _ensure_bullet_formatting(self, paragraph: Any, reference_paragraph: Any = None) -> None:
        """Copy complete bullet formatting from reference paragraph to target.

        Simpler approach: directly copy the entire pPr (paragraph properties) element
        from the template paragraph. This ensures ALL formatting is preserved.

        Args:
            paragraph: PowerPoint paragraph to add bullet formatting to
            reference_paragraph: Template paragraph to copy bullet formatting from
        """
        if not reference_paragraph:
            logger.warning("No reference paragraph provided for bullet formatting")
            return

        try:
            from copy import deepcopy
        except ImportError:
            logger.warning("deepcopy not available, cannot copy bullet formatting")
            return

        try:
            # Get XML elements
            ref_elem = reference_paragraph._element
            target_elem = paragraph._element

            # Find pPr (paragraph properties) in reference paragraph
            ref_pPr = None
            for child in ref_elem:
                if child.tag.endswith('}pPr'):
                    ref_pPr = child
                    break

            if ref_pPr is None:
                logger.warning("Reference paragraph has no pPr element")
                return

            # Remove any existing pPr in target paragraph
            for child in list(target_elem):
                if child.tag.endswith('}pPr'):
                    target_elem.remove(child)

            # Deep copy entire pPr element from reference to target
            # This preserves all bullet properties: character, font, color, size, etc.
            new_pPr = deepcopy(ref_pPr)
            target_elem.insert(0, new_pPr)

            logger.debug("Successfully copied complete bullet formatting from reference paragraph")

        except Exception as e:
            logger.error(f"Failed to copy bullet formatting: {e}", exc_info=True)

    def _set_bullet_list(self, text_frame: Any, bullet_items: List[str]) -> None:
        """Set bullet list content while preserving template formatting.

        Uses the template's first paragraph as a reference to copy bullet formatting
        to all dynamically added paragraphs. This ensures consistent formatting.

        Args:
            text_frame: PowerPoint text frame to populate
            bullet_items: List of bullet point strings (may contain **bold** markdown)
        """
        # Capture original formatting from first paragraph's first run
        original_font_name = None
        original_font_size = None
        original_font_color = None
        reference_paragraph = None

        if text_frame.paragraphs and text_frame.paragraphs[0].runs:
            # Keep reference to first paragraph for bullet formatting
            reference_paragraph = text_frame.paragraphs[0]
            first_run = reference_paragraph.runs[0]
            original_font_name = first_run.font.name
            original_font_size = first_run.font.size
            if hasattr(first_run.font.color, 'rgb'):
                try:
                    original_font_color = first_run.font.color.rgb
                except:
                    pass  # Color might be theme-based, skip

        # Build list of paragraphs, reusing existing ones where possible
        paragraphs = []
        for i in range(len(bullet_items)):
            if i < len(text_frame.paragraphs):
                # Reuse existing paragraph (already has bullet formatting from template)
                p = text_frame.paragraphs[i]
                # Clear existing runs
                for run in p.runs[:]:
                    run.text = ""
            else:
                # Add new paragraph for additional bullets
                p = text_frame.add_paragraph()
                p.level = 0  # Set to first-level bullet
                # Copy complete bullet formatting from template's first paragraph
                self._ensure_bullet_formatting(p, reference_paragraph)

            paragraphs.append(p)

        # Populate each paragraph with bullet content
        for p, bullet_text in zip(paragraphs, bullet_items):
            bullet_text = bullet_text.strip()

            # Parse **bold** markdown formatting
            parts = re.split(r'(\*\*.*?\*\*)', bullet_text)

            for part in parts:
                if not part:
                    continue

                if part.startswith('**') and part.endswith('**'):
                    # Bold text
                    run = p.add_run()
                    run.text = part[2:-2]
                    run.font.bold = True
                    # Apply captured formatting
                    if original_font_name:
                        run.font.name = original_font_name
                    if original_font_size:
                        run.font.size = original_font_size
                    if original_font_color:
                        run.font.color.rgb = original_font_color
                else:
                    # Regular text
                    run = p.add_run()
                    run.text = part
                    # Apply captured formatting
                    if original_font_name:
                        run.font.name = original_font_name
                    if original_font_size:
                        run.font.size = original_font_size
                    if original_font_color:
                        run.font.color.rgb = original_font_color

    def _populate_template(
        self,
        slide: Any,
        headline_statement: str,
        process_name: str,
        nhs_org_name: str,
        process_category: str,
        systems_used: List[str],
        manual_process: str,
        challenges: str,
        rpa_solution: str,
        benefits: List[str]
    ) -> None:
        """Replace placeholders in the template with provided content."""

        # Define placeholder mappings (based on template analysis)
        # Note: benefits handled separately with _set_bullet_list()
        replacements = {
            "[Headline Statement Placeholder]": headline_statement,
            "Process Name: [Waiting List Validation]": f"Process Name: {process_name}",
            "[NHS Org Name Placeholder]": nhs_org_name,
            "[Process Category Placeholder]": process_category,
            "Systems Used: [PAS | RADAR | Patient Hub]": f"Systems Used: {' | '.join(systems_used)}",
            "[Manual processes placeholder]": manual_process,
            "[Challenges placeholder]": challenges,
            "[RPA Solution placeholder]": rpa_solution,
        }

        # Iterate through shapes and replace text
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue

            text_frame = shape.text_frame

            # Check if this shape contains a placeholder we need to replace
            for placeholder, replacement in replacements.items():
                if placeholder in text_frame.text:
                    # Handle markdown formatting
                    if any(x in replacement for x in ["**", "\n-", "\n*"]) or replacement.strip().startswith(('- ', '* ')):
                        self._set_formatted_text(text_frame, replacement)
                    else:
                        # Simple text replacement - preserve formatting by using paragraph-level replacement
                        replaced = False
                        for paragraph in text_frame.paragraphs:
                            if placeholder in paragraph.text:
                                # Capture formatting from first run
                                font_name = None
                                font_size = None
                                font_color = None
                                font_bold = None

                                if paragraph.runs:
                                    first_run = paragraph.runs[0]
                                    font_name = first_run.font.name
                                    font_size = first_run.font.size
                                    font_bold = first_run.font.bold
                                    if hasattr(first_run.font.color, 'rgb'):
                                        try:
                                            font_color = first_run.font.color.rgb
                                        except:
                                            pass

                                # Replace in full paragraph text
                                new_text = paragraph.text.replace(placeholder, replacement)

                                # Clear all runs
                                for run in paragraph.runs[:]:
                                    run.text = ""

                                # Create or reuse first run
                                if paragraph.runs:
                                    run = paragraph.runs[0]
                                else:
                                    run = paragraph.add_run()

                                run.text = new_text

                                # Restore formatting
                                if font_name:
                                    run.font.name = font_name
                                if font_size:
                                    run.font.size = font_size
                                if font_bold is not None:
                                    run.font.bold = font_bold
                                if font_color:
                                    run.font.color.rgb = font_color

                                replaced = True
                                break

                        if not replaced:
                            # Fallback if placeholder not found in any paragraph
                            text_frame.text = text_frame.text.replace(placeholder, replacement)

                    # Auto-adjust font size if text is slightly over limit
                    self._auto_adjust_font_size(shape, text_frame, replacement)
                    break

            # Handle benefits placeholder separately (requires bullet list formatting)
            if "[Benefits placeholder]" in text_frame.text:
                self._set_bullet_list(text_frame, benefits)
                # Auto-adjust font size if text is slightly over limit
                # Estimate total characters for sizing
                total_chars = sum(len(b) for b in benefits)
                self._auto_adjust_font_size(shape, text_frame, ' '.join(benefits))

    def _set_formatted_text(self, text_frame: Any, markdown_text: str) -> None:
        """Convert markdown to PowerPoint formatted text while preserving bullet formatting."""
        # Save reference to first paragraph BEFORE clearing (for bullet formatting)
        reference_paragraph = None
        original_font_name = None
        original_font_size = None
        original_font_color = None

        if text_frame.paragraphs and text_frame.paragraphs[0].runs:
            reference_paragraph = text_frame.paragraphs[0]
            first_run = reference_paragraph.runs[0]
            original_font_name = first_run.font.name
            original_font_size = first_run.font.size
            if hasattr(first_run.font.color, 'rgb'):
                try:
                    original_font_color = first_run.font.color.rgb
                except:
                    pass  # Color might be theme-based, skip

        # Clear existing paragraphs
        text_frame.clear()

        # Split by lines
        lines = markdown_text.split('\n')

        # Track if we're on the first content paragraph
        first_paragraph = True

        for line in lines:
            line = line.strip()
            if not line:
                # Empty line - add blank paragraph
                text_frame.add_paragraph()
                first_paragraph = False
                continue

            # Check if it's a bullet point
            is_bullet = line.startswith('- ') or line.startswith('* ')
            if is_bullet:
                line = line[2:].strip()  # Remove bullet marker

            # Create or reuse paragraph
            # After clear(), one empty paragraph remains - reuse it for first content line
            if first_paragraph and text_frame.paragraphs:
                p = text_frame.paragraphs[0]
                first_paragraph = False
            else:
                p = text_frame.add_paragraph()

            if is_bullet:
                p.level = 0  # First level bullet
                # Copy bullet formatting from saved reference paragraph
                self._ensure_bullet_formatting(p, reference_paragraph)

            # Process bold formatting
            parts = re.split(r'(\*\*.*?\*\*)', line)
            for part in parts:
                if part.startswith('**') and part.endswith('**'):
                    # Bold text
                    run = p.add_run()
                    run.text = part[2:-2]
                    run.font.bold = True
                    # Apply captured formatting
                    if original_font_name:
                        run.font.name = original_font_name
                    if original_font_size:
                        run.font.size = original_font_size
                    if original_font_color:
                        run.font.color.rgb = original_font_color
                elif part:
                    # Regular text
                    run = p.add_run()
                    run.text = part
                    # Apply captured formatting
                    if original_font_name:
                        run.font.name = original_font_name
                    if original_font_size:
                        run.font.size = original_font_size
                    if original_font_color:
                        run.font.color.rgb = original_font_color

    def _auto_adjust_font_size(self, shape: Any, text_frame: Any, text: str) -> None:
        """Automatically adjust font size if text is slightly over limit."""
        # Get current dimensions
        width_inches = shape.width.inches
        height_inches = shape.height.inches

        # Simple heuristic: if text is long relative to box size, reduce font slightly
        char_count = len(text)
        area = width_inches * height_inches
        density = char_count / area if area > 0 else 0

        # If density is high (>200 chars per square inch), reduce font size
        if density > 200:
            reduction_factor = min(0.85, 200 / density)  # Max 15% reduction

            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size:
                        current_size = run.font.size.pt
                        new_size = max(7, current_size * reduction_factor)  # Min 7pt
                        run.font.size = Pt(new_size)

    async def _upload_to_storage(
        self,
        user_id: str,
        filename: str,
        file_bytes: bytes
    ) -> str:
        """Upload PowerPoint file to Supabase storage and return signed URL."""
        try:
            from supabase import create_client

            # Check if Supabase is configured
            if not settings.supabase_url or not settings.supabase_service_key:
                raise ToolExecutionError(
                    "generate_process_one_pager",
                    "Supabase storage not configured. Contact administrator."
                )

            # Create Supabase client
            supabase = create_client(
                settings.supabase_url,
                settings.supabase_service_key
            )

            # Upload file to user's folder
            storage_path = f"{user_id}/generated_ppts/{filename}"

            supabase.storage.from_("process-one-pagers").upload(
                path=storage_path,
                file=file_bytes,
                file_options={
                    "content-type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    "cache-control": "3600",
                    "upsert": "true"  # Overwrite if exists
                }
            )

            # Generate signed URL (1 hour expiry)
            response = supabase.storage.from_("process-one-pagers").create_signed_url(
                path=storage_path,
                expires_in=3600
            )

            signed_url = response.get("signedURL")

            if not signed_url:
                raise ToolExecutionError(
                    "generate_process_one_pager",
                    "Failed to generate download URL from storage"
                )

            # Fix URL for development environment
            if os.getenv("ENVIRONMENT", "development") == "development":
                signed_url = signed_url.replace("kong:8000", "localhost:8000")

            return signed_url

        except Exception as e:
            logger.error(f"Failed to upload to storage: {e}", exc_info=True)
            raise ToolExecutionError(
                "generate_process_one_pager",
                f"Failed to upload file to storage: {str(e)}"
            )
