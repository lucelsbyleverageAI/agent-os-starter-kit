"""Outlook email tools for the MCP server."""

from typing import Any, Dict, List, Optional
from datetime import datetime, timedelta, timezone
import tempfile
import os
from pathlib import Path

import mistune

from ..base import CustomTool, ToolParameter
from ...utils.logging import get_logger
from ...utils.exceptions import ToolExecutionError

from .base import get_outlook_client, handle_outlook_error

logger = get_logger(__name__)

# Import lightweight document parsing libraries
try:
    from pypdf import PdfReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    logger.warning("pypdf not available. PDF extraction will be disabled.")

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not available. DOCX extraction will be disabled.")

try:
    from openpyxl import load_workbook
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False
    logger.warning("openpyxl not available. XLSX extraction will be disabled.")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available. PPTX extraction will be disabled.")

try:
    from bs4 import BeautifulSoup
    HTML_AVAILABLE = True
except ImportError:
    HTML_AVAILABLE = False
    logger.warning("BeautifulSoup not available. HTML extraction will be disabled.")


def markdown_to_html(markdown_text: str) -> str:
    """Convert markdown to HTML for email body using mistune."""
    return mistune.html(markdown_text)


def format_recipients_list(recipients: List[Dict]) -> str:
    """Format recipient list for display."""
    if not recipients:
        return "None"

    emails = [r.get("emailAddress", {}).get("address", "Unknown") for r in recipients]
    return ", ".join(emails)


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF using pypdf."""
    if not PDF_AVAILABLE:
        return "PDF parsing library not available"

    try:
        reader = PdfReader(file_path)
        text_parts = []

        for page_num, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            if text.strip():
                text_parts.append(f"## Page {page_num}\n\n{text.strip()}\n")

        return "\n".join(text_parts) if text_parts else "No text content found in PDF"
    except Exception as e:
        return f"Error extracting PDF: {str(e)}"


def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX using python-docx."""
    if not DOCX_AVAILABLE:
        return "DOCX parsing library not available"

    try:
        doc = DocxDocument(file_path)
        text_parts = []

        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                # Check if it's a heading
                if para.style.name.startswith('Heading'):
                    level = para.style.name.replace('Heading ', '')
                    if level.isdigit():
                        text_parts.append(f"{'#' * int(level)} {para.text.strip()}\n")
                    else:
                        text_parts.append(f"## {para.text.strip()}\n")
                else:
                    text_parts.append(f"{para.text.strip()}\n")

        # Extract tables
        for table_num, table in enumerate(doc.tables, 1):
            text_parts.append(f"\n### Table {table_num}\n")
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                text_parts.append("| " + " | ".join(cells) + " |")
            text_parts.append("")

        return "\n".join(text_parts) if text_parts else "No text content found in DOCX"
    except Exception as e:
        return f"Error extracting DOCX: {str(e)}"


def extract_text_from_xlsx(file_path: str) -> str:
    """Extract text from XLSX using openpyxl."""
    if not XLSX_AVAILABLE:
        return "XLSX parsing library not available"

    try:
        wb = load_workbook(file_path, data_only=True)
        text_parts = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_parts.append(f"## Sheet: {sheet_name}\n")

            # Get all rows with data
            for row in sheet.iter_rows(values_only=True):
                # Skip empty rows
                if any(cell is not None for cell in row):
                    cells = [str(cell) if cell is not None else "" for cell in row]
                    text_parts.append("| " + " | ".join(cells) + " |")

            text_parts.append("")

        return "\n".join(text_parts) if text_parts else "No data found in XLSX"
    except Exception as e:
        return f"Error extracting XLSX: {str(e)}"


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX using python-pptx."""
    if not PPTX_AVAILABLE:
        return "PPTX parsing library not available"

    try:
        prs = Presentation(file_path)
        text_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts.append(f"## Slide {slide_num}\n")

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(f"{shape.text.strip()}\n")

            # Extract notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                text_parts.append(f"\n**Notes:** {slide.notes_slide.notes_text_frame.text.strip()}\n")

            text_parts.append("")

        return "\n".join(text_parts) if text_parts else "No text content found in PPTX"
    except Exception as e:
        return f"Error extracting PPTX: {str(e)}"


def extract_text_from_html(file_path: str) -> str:
    """Extract text from HTML using BeautifulSoup."""
    if not HTML_AVAILABLE:
        return "HTML parsing library not available"

    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            html_content = f.read()

        soup = BeautifulSoup(html_content, 'lxml')

        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()

        # Get text
        text = soup.get_text()

        # Clean up whitespace
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = '\n'.join(chunk for chunk in chunks if chunk)

        return text if text.strip() else "No text content found in HTML"
    except Exception as e:
        return f"Error extracting HTML: {str(e)}"


def truncate_text_with_metadata(text: str, start_word: int = 0, max_words: int = 2000) -> Dict[str, Any]:
    """Truncate text to specified word range and return with metadata.

    Args:
        text: The full text to truncate
        start_word: Starting word index (0-based)
        max_words: Maximum number of words to return

    Returns:
        Dictionary containing:
        - text: The truncated text
        - total_words: Total word count in full text
        - start_word: Starting word index
        - end_word: Ending word index (exclusive)
        - is_truncated: Whether text was truncated
        - has_more: Whether there are more words available
    """
    words = text.split()
    total_words = len(words)

    # Calculate end word
    end_word = min(start_word + max_words, total_words)

    # Extract requested range
    truncated_words = words[start_word:end_word]
    truncated_text = " ".join(truncated_words)

    return {
        "text": truncated_text,
        "total_words": total_words,
        "start_word": start_word,
        "end_word": end_word,
        "is_truncated": end_word < total_words,
        "has_more": end_word < total_words,
    }


async def extract_attachment_text(
    client,
    user_email: str,
    message_id: str,
    attachment_id: str,
    attachment_name: str,
    max_words: int = 2000,
) -> str:
    """Download and extract text from an attachment.

    Args:
        client: Outlook client instance
        user_email: User email address
        message_id: Email message ID
        attachment_id: Attachment ID
        attachment_name: Attachment filename
        max_words: Maximum words to extract (default 2000)

    Returns:
        Extracted text with metadata (truncated if necessary)
    """
    try:
        # Download binary content
        binary_content = await client.download_attachment_binary(
            user_email=user_email,
            message_id=message_id,
            attachment_id=attachment_id,
        )

        # Get file extension
        file_extension = Path(attachment_name).suffix.lower()

        # Save to temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as temp_file:
            temp_file.write(binary_content)
            temp_file_path = temp_file.name

        try:
            # Extract text based on file type
            extracted_text = ""

            if file_extension == ".pdf":
                extracted_text = extract_text_from_pdf(temp_file_path)
            elif file_extension == ".docx":
                extracted_text = extract_text_from_docx(temp_file_path)
            elif file_extension in [".xlsx", ".xls"]:
                extracted_text = extract_text_from_xlsx(temp_file_path)
            elif file_extension in [".pptx", ".ppt"]:
                extracted_text = extract_text_from_pptx(temp_file_path)
            elif file_extension in [".html", ".htm"]:
                extracted_text = extract_text_from_html(temp_file_path)
            elif file_extension in [".txt", ".md", ".csv", ".eml"]:
                with open(temp_file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    extracted_text = f.read()
            else:
                return f"**{attachment_name}**\n\nUnsupported format: {file_extension}\n"

            # Truncate if necessary
            if extracted_text:
                truncated = truncate_text_with_metadata(extracted_text, 0, max_words)

                result_parts = [f"**{attachment_name}**"]

                if truncated["is_truncated"]:
                    result_parts.append(
                        f"*Showing {truncated['start_word']}-{truncated['end_word']} of {truncated['total_words']} words. "
                        f"Use download_outlook_attachment with start_word={truncated['end_word']} to continue.*\n"
                    )
                else:
                    result_parts.append(f"*Complete text ({truncated['total_words']} words)*\n")

                result_parts.append(truncated["text"])
                return "\n".join(result_parts)
            else:
                return f"**{attachment_name}**\n\nNo text content extracted.\n"

        finally:
            # Clean up temp file
            os.unlink(temp_file_path)

    except Exception as e:
        logger.error(f"Error extracting attachment {attachment_name}: {str(e)}")
        return f"**{attachment_name}**\n\nError extracting text: {str(e)}\n"


class ListEmailsTool(CustomTool):
    """List user's emails with optional filtering across all folders."""

    toolkit_name = "outlook"
    toolkit_display_name = "E18 Outlook"

    @property
    def name(self) -> str:
        return "list_emails"

    @property
    def description(self) -> str:
        return (
            "List your Outlook emails with optional filtering. Searches across ALL mail folders by default. "
            "Returns email metadata including subject, senders, recipients, and preview. "
            "\n\n"
            "FILTERS (can be combined):\n"
            "- Date range: Use from_date/to_date (ISO 8601 format) or time_range shortcut (day/week/month/year)\n"
            "- Sender: Filter by sender email address (single email only)\n"
            "- Subject/body keywords: Search in subject line or body content\n"
            "- Folder: Specify folder ID to search specific folder (use list_mail_folders to get IDs)\n"
            "- ALL filters can now be combined (e.g., date range + sender + keywords)\n"
            "\n"
            "PAGINATION:\n"
            "- Use 'next_link' parameter to get the next page of results\n"
            "- The 'next_link' URL is returned in the response when more results exist\n"
            "- Cursor-based pagination is fast and efficient (no performance degradation)\n"
            "- Maximum 1,000 results per search query (use date ranges to narrow scope if needed)\n"
            "\n"
            "PERFORMANCE NOTES:\n"
            "- Uses optimized KQL-based search for fast filtering on large mailboxes\n"
            "- Results are sorted by date (most recent first)\n"
            "- Complex filter combinations no longer cause timeouts\n"
            "\n"
            "Use message IDs with read_email tool to get full content."
        )

    def get_parameters(self) -> List[ToolParameter]:
        return [
            ToolParameter(
                name="from_date",
                type="string",
                description="Start date in ISO 8601 format (e.g., '2024-01-01T00:00:00Z'). Only emails received on or after this date.",
                required=False,
            ),
            ToolParameter(
                name="to_date",
                type="string",
                description="End date in ISO 8601 format. Only emails received on or before this date.",
                required=False,
            ),
            ToolParameter(
                name="time_range",
                type="string",
                description="Convenient time range (e.g., 'day', 'week', 'month'). Overrides from_date/to_date if provided.",
                required=False,
                enum=["day", "week", "month", "year"],
            ),
            ToolParameter(
                name="sender",
                type="string",
                description="Filter by sender email address (e.g., 'boss@company.com'). Single email address only.",
                required=False,
            ),
            ToolParameter(
                name="subject_keyword",
                type="string",
                description="Search keyword in email subject. Can be combined with date/sender filters.",
                required=False,
            ),
            ToolParameter(
                name="body_keyword",
                type="string",
                description="Search keyword in email body content. Can be combined with date/sender filters.",
                required=False,
            ),
            ToolParameter(
                name="folder_id",
                type="string",
                description="Folder ID to search (leave empty to search all folders). Use list_mail_folders to get folder IDs.",
                required=False,
            ),
            ToolParameter(
                name="limit",
                type="integer",
                description="Maximum number of emails to return (max 1000). Defaults to 20.",
                required=False,
                default=20,
            ),
            ToolParameter(
                name="next_link",
                type="string",
                description="URL from previous response's 'next_link' field to get the next page. Leave empty for first page.",
                required=False,
            ),
        ]

    async def _execute_impl(self, user_id: str, **kwargs: Any) -> str:
        """Execute list emails tool."""
        try:
            # Get user email from context - REQUIRED for scoping
            user_email = kwargs.get("_context_user_email")

            if not user_email:
                raise ToolExecutionError(
                    "list_emails",
                    "User email is required for data scoping. Please ensure you are properly authenticated."
                )

            # Parse time range
            time_range = kwargs.get("time_range")
            from_date = kwargs.get("from_date")
            to_date = kwargs.get("to_date")

            if time_range:
                now = datetime.now(timezone.utc)
                if time_range == "day":
                    from_date = (now - timedelta(days=1)).isoformat()
                elif time_range == "week":
                    from_date = (now - timedelta(weeks=1)).isoformat()
                elif time_range == "month":
                    from_date = (now - timedelta(days=30)).isoformat()
                elif time_range == "year":
                    from_date = (now - timedelta(days=365)).isoformat()
                to_date = now.isoformat()

            # Get parameters
            sender = kwargs.get("sender")
            subject_keyword = kwargs.get("subject_keyword")
            body_keyword = kwargs.get("body_keyword")
            folder_id = kwargs.get("folder_id")
            limit = kwargs.get("limit", 20)
            next_link = kwargs.get("next_link")

            # Query Outlook
            client = get_outlook_client()
            response = await client.list_messages(
                user_email=user_email,
                from_date=from_date,
                to_date=to_date,
                sender=sender,
                subject_keyword=subject_keyword,
                body_keyword=body_keyword,
                folder_id=folder_id,
                limit=limit,
                next_link=next_link,
            )

            messages = response.get("value", [])

            if not messages:
                return "*No emails found matching your criteria.*"

            # Format as markdown
            markdown_parts = [f"# Outlook Emails ({len(messages)} found)\n"]

            for msg in messages:
                message_id = msg.get("id", "")
                subject = msg.get("subject", "No Subject")
                from_info = msg.get("from", {}).get("emailAddress", {})
                from_email = from_info.get("address", "Unknown")
                from_name = from_info.get("name", from_email)
                received = msg.get("receivedDateTime", "")
                preview = msg.get("bodyPreview", "")[:200]  # First 200 chars
                has_attachments = msg.get("hasAttachments", False)
                is_read = msg.get("isRead", False)
                importance = msg.get("importance", "normal")

                # Format received date
                try:
                    received_dt = datetime.fromisoformat(received.replace("Z", "+00:00"))
                    received_str = received_dt.strftime("%Y-%m-%d %H:%M")
                except:
                    received_str = received

                markdown_parts.append(f"## {subject}")
                markdown_parts.append(f"**Message ID:** `{message_id}`")
                markdown_parts.append(f"**From:** {from_name} <{from_email}>")
                markdown_parts.append(f"**Received:** {received_str}")

                to_recipients = format_recipients_list(msg.get("toRecipients", []))
                cc_recipients = format_recipients_list(msg.get("ccRecipients", []))
                markdown_parts.append(f"**To:** {to_recipients}")
                if cc_recipients != "None":
                    markdown_parts.append(f"**CC:** {cc_recipients}")

                status_flags = []
                if not is_read:
                    status_flags.append("Unread")
                if has_attachments:
                    status_flags.append("Has Attachments")
                if importance != "normal":
                    status_flags.append(f"Importance: {importance}")

                if status_flags:
                    markdown_parts.append(f"**Status:** {', '.join(status_flags)}")

                if preview:
                    markdown_parts.append(f"\n**Preview:** {preview}...")

                markdown_parts.append("")

            # Add pagination info with next_link
            odata_next_link = response.get("@odata.nextLink")
            if odata_next_link:
                markdown_parts.append("\n---\n")
                markdown_parts.append("## Pagination")
                markdown_parts.append(f"\n**More results available!** To get the next page, use:")
                markdown_parts.append(f"\n```")
                markdown_parts.append(f"next_link: {odata_next_link}")
                markdown_parts.append(f"```")
                markdown_parts.append(f"\n*Copy the above URL and paste it into the 'next_link' parameter.*")

            return "\n".join(markdown_parts)

        except Exception as e:
            error_msg = handle_outlook_error(e)
            logger.error(f"Error in list_emails: {error_msg}")
            raise ToolExecutionError("list_emails", error_msg)


class ListMailFoldersTool(CustomTool):
    """List all mail folders for the user."""

    toolkit_name = "outlook"
    toolkit_display_name = "E18 Outlook"

    @property
    def name(self) -> str:
        return "list_mail_folders"

    @property
    def description(self) -> str:
        return (
            "List all Outlook mail folders for your account. "
            "Returns folder names, IDs, and message counts. "
            "\n\n"
            "Use folder IDs with the list_emails tool to filter emails by folder. "
            "\n\n"
            "Common folders include: Inbox, Sent Items, Drafts, Deleted Items, Archive, and custom folders."
        )

    def get_parameters(self) -> List[ToolParameter]:
        return []

    async def _execute_impl(self, user_id: str, **kwargs: Any) -> str:
        """Execute list mail folders tool."""
        try:
            user_email = kwargs.get("_context_user_email")

            if not user_email:
                raise ToolExecutionError(
                    "list_mail_folders",
                    "User email is required for data scoping."
                )

            # Get folders
            client = get_outlook_client()
            response = await client.list_mail_folders(user_email=user_email)

            folders = response.get("value", [])

            if not folders:
                return "*No mail folders found.*"

            # Format as markdown
            markdown_parts = [f"# Outlook Mail Folders ({len(folders)} found)\n"]

            for folder in folders:
                folder_id = folder.get("id", "")
                display_name = folder.get("displayName", "Unknown")
                total_items = folder.get("totalItemCount", 0)
                unread_items = folder.get("unreadItemCount", 0)

                markdown_parts.append(f"## {display_name}")
                markdown_parts.append(f"**Folder ID:** `{folder_id}`")
                markdown_parts.append(f"**Total Messages:** {total_items}")
                markdown_parts.append(f"**Unread:** {unread_items}")
                markdown_parts.append("")

            return "\n".join(markdown_parts)

        except Exception as e:
            error_msg = handle_outlook_error(e)
            logger.error(f"Error in list_mail_folders: {error_msg}")
            raise ToolExecutionError("list_mail_folders", error_msg)


class ReadEmailTool(CustomTool):
    """Read full email content including body and attachments."""

    toolkit_name = "outlook"
    toolkit_display_name = "E18 Outlook"

    @property
    def name(self) -> str:
        return "read_email"

    @property
    def description(self) -> str:
        return (
            "Read the full content of an Outlook email. "
            "Returns complete email details including full body text, all recipients, and attachment metadata. "
            "\n\n"
            "ATTACHMENT HANDLING:\n"
            "Option 1 - Auto-download all attachments (simple):\n"
            "  Set include_attachments=true to automatically download and extract text from ALL attachments.\n"
            "  Each attachment is truncated to 2000 words by default.\n"
            "\n"
            "Option 2 - Download specific attachments (precise control):\n"
            "  Use download_outlook_attachment tool with message_id + attachment_name.\n"
            "  Supports pagination for large documents.\n"
            "\n"
            "SUPPORTED FORMATS:\n"
            "- PDF, DOCX, XLSX, PPTX, HTML, CSV, TXT, EML\n"
            "\n\n"
            "Use the message ID from list_emails tool."
        )

    def get_parameters(self) -> List[ToolParameter]:
        return [
            ToolParameter(
                name="message_id",
                type="string",
                description="The message ID from list_emails.",
                required=True,
            ),
            ToolParameter(
                name="body_format",
                type="string",
                description="Body format preference: 'text' or 'html'. Defaults to 'text'.",
                required=False,
                enum=["text", "html"],
                default="text",
            ),
            ToolParameter(
                name="include_attachments",
                type="boolean",
                description="If true, automatically downloads and extracts text from all attachments. Each attachment is truncated to 2000 words. Default: false.",
                required=False,
                default=False,
            ),
        ]

    async def _execute_impl(self, user_id: str, **kwargs: Any) -> str:
        """Execute read email tool."""
        try:
            user_email = kwargs.get("_context_user_email")

            if not user_email:
                raise ToolExecutionError(
                    "read_email",
                    "User email is required for data scoping."
                )

            message_id = kwargs.get("message_id")
            body_format = kwargs.get("body_format", "text")
            include_attachments = kwargs.get("include_attachments", False)

            if not message_id:
                raise ToolExecutionError("read_email", "message_id is required")

            # Get message
            client = get_outlook_client()
            msg = await client.get_message(
                user_email=user_email,
                message_id=message_id,
                body_format=body_format,
            )

            # Extract details
            subject = msg.get("subject", "No Subject")
            from_info = msg.get("from", {}).get("emailAddress", {})
            from_email = from_info.get("address", "Unknown")
            from_name = from_info.get("name", from_email)
            received = msg.get("receivedDateTime", "")
            body_obj = msg.get("body", {})
            body_content = body_obj.get("content", "")

            # Format markdown
            markdown_parts = [f"# Email: {subject}\n"]
            markdown_parts.append(f"**From:** {from_name} <{from_email}>")
            markdown_parts.append(f"**To:** {format_recipients_list(msg.get('toRecipients', []))}")

            cc = format_recipients_list(msg.get("ccRecipients", []))
            if cc != "None":
                markdown_parts.append(f"**CC:** {cc}")

            bcc = format_recipients_list(msg.get("bccRecipients", []))
            if bcc != "None":
                markdown_parts.append(f"**BCC:** {bcc}")

            try:
                received_dt = datetime.fromisoformat(received.replace("Z", "+00:00"))
                received_str = received_dt.strftime("%Y-%m-%d %H:%M:%S")
            except:
                received_str = received

            markdown_parts.append(f"**Received:** {received_str}")
            markdown_parts.append(f"**Message ID:** `{message_id}`\n")
            markdown_parts.append("---\n")

            # Body
            markdown_parts.append("## Body\n")
            markdown_parts.append(body_content)
            markdown_parts.append("")

            # Attachments (fetch metadata or full text)
            if msg.get("hasAttachments"):
                try:
                    attachments_response = await client.get_message_attachments(
                        user_email=user_email,
                        message_id=message_id,
                    )
                    attachments = attachments_response.get("value", [])

                    if attachments:
                        markdown_parts.append("\n## Attachments\n")

                        if include_attachments:
                            # Download and extract text from all attachments
                            for att in attachments:
                                att_name = att.get("name", "Unknown")
                                att_id = att.get("id", "")
                                att_size = att.get("size", 0)

                                # Skip if too large
                                if att_size > 3 * 1024 * 1024:
                                    markdown_parts.append(
                                        f"**{att_name}**\n"
                                        f"*Skipped: File size ({att_size / (1024 * 1024):.1f} MB) exceeds 3 MB limit*\n"
                                    )
                                    continue

                                # Extract text
                                extracted_text = await extract_attachment_text(
                                    client=client,
                                    user_email=user_email,
                                    message_id=message_id,
                                    attachment_id=att_id,
                                    attachment_name=att_name,
                                    max_words=2000,
                                )
                                markdown_parts.append(extracted_text)
                        else:
                            # Just show metadata
                            for att in attachments:
                                att_name = att.get("name", "Unknown")
                                att_type = att.get("contentType", "Unknown")
                                att_size = att.get("size", 0)
                                is_inline = att.get("isInline", False)

                                # Format size
                                if att_size < 1024:
                                    size_str = f"{att_size} bytes"
                                elif att_size < 1024 * 1024:
                                    size_str = f"{att_size / 1024:.1f} KB"
                                else:
                                    size_str = f"{att_size / (1024 * 1024):.1f} MB"

                                # Check if size exceeds 3 MB limit
                                size_warning = ""
                                if att_size > 3 * 1024 * 1024:
                                    size_warning = " ⚠️ (Exceeds 3 MB download limit)"

                                markdown_parts.append(f"### {att_name}")
                                markdown_parts.append(f"- **Type:** {att_type}")
                                markdown_parts.append(f"- **Size:** {size_str}{size_warning}")
                                if is_inline:
                                    markdown_parts.append(f"- **Inline:** Yes (embedded in email)")
                                markdown_parts.append("")

                            markdown_parts.append(
                                "*To extract text from specific attachments, use the `download_outlook_attachment` tool with message_id + attachment_name.*"
                            )
                except Exception as att_error:
                    logger.warning(f"Error fetching attachments: {att_error}")
                    markdown_parts.append(
                        "\n## Attachments\n*This email has attachments, but metadata could not be retrieved.*"
                    )

            return "\n".join(markdown_parts)

        except Exception as e:
            error_msg = handle_outlook_error(e)
            logger.error(f"Error in read_email: {error_msg}")
            raise ToolExecutionError("read_email", error_msg)


class DownloadOutlookAttachmentTool(CustomTool):
    """Download and extract text from email attachments."""

    toolkit_name = "outlook"
    toolkit_display_name = "E18 Outlook"

    @property
    def name(self) -> str:
        return "download_outlook_attachment"

    @property
    def description(self) -> str:
        return (
            "Download an Outlook email attachment and extract its text content with pagination support. "
            "\n\n"
            "SUPPORTED FORMATS:\n"
            "- Documents: PDF, DOCX, XLSX, PPTX\n"
            "- Web: HTML\n"
            "- Text: TXT, Markdown, CSV, EML\n"
            "\n\n"
            "PAGINATION:\n"
            "- Default: Returns first 2000 words of extracted text\n"
            "- For longer documents, use start_word parameter to navigate\n"
            "- Example: start_word=2000 returns words 2000-4000\n"
            "\n\n"
            "SIZE LIMITS:\n"
            "- Maximum attachment size: 3 MB (Microsoft Graph API limit)\n"
            "\n\n"
            "USAGE EXAMPLES:\n"
            "1. Download attachment from email:\n"
            "   message_id='ABC123', attachment_name='report.pdf'\n"
            "\n"
            "2. Navigate long document:\n"
            "   message_id='ABC123', attachment_name='report.pdf', start_word=2000\n"
            "\n"
            "Returns extracted text with metadata showing total words and current range."
        )

    def get_parameters(self) -> List[ToolParameter]:
        return [
            ToolParameter(
                name="message_id",
                type="string",
                description="The email message ID from read_email or list_emails.",
                required=True,
            ),
            ToolParameter(
                name="attachment_name",
                type="string",
                description="The exact filename of the attachment (e.g., 'report.pdf', 'data.xlsx').",
                required=True,
            ),
            ToolParameter(
                name="start_word",
                type="number",
                description="Starting word index for pagination (0-based). Default: 0.",
                required=False,
                default=0,
            ),
            ToolParameter(
                name="max_words",
                type="number",
                description="Maximum words to return per request. Default: 2000.",
                required=False,
                default=2000,
            ),
        ]

    async def _execute_impl(self, user_id: str, **kwargs: Any) -> str:
        """Execute download attachment tool."""
        try:
            user_email = kwargs.get("_context_user_email")

            if not user_email:
                raise ToolExecutionError(
                    "download_outlook_attachment",
                    "User email is required for data scoping."
                )

            message_id = kwargs.get("message_id")
            attachment_name = kwargs.get("attachment_name")
            start_word = int(kwargs.get("start_word", 0))
            max_words = int(kwargs.get("max_words", 2000))

            if not message_id:
                raise ToolExecutionError(
                    "download_outlook_attachment",
                    "message_id is required"
                )

            if not attachment_name:
                raise ToolExecutionError(
                    "download_outlook_attachment",
                    "attachment_name is required"
                )

            # Get all attachments and find by name
            client = get_outlook_client()
            attachments_response = await client.get_message_attachments(
                user_email=user_email,
                message_id=message_id,
            )

            attachments = attachments_response.get("value", [])
            attachment = next((a for a in attachments if a.get("name") == attachment_name), None)

            if not attachment:
                available_names = [a.get("name", "Unknown") for a in attachments]
                raise ToolExecutionError(
                    "download_outlook_attachment",
                    f"Attachment '{attachment_name}' not found. Available attachments: {', '.join(available_names)}"
                )

            attachment_id = attachment.get("id", "")

            att_name = attachment.get("name", "Unknown")
            att_type = attachment.get("contentType", "Unknown")
            att_size = attachment.get("size", 0)

            # Format size
            if att_size < 1024:
                size_str = f"{att_size} bytes"
            elif att_size < 1024 * 1024:
                size_str = f"{att_size / 1024:.1f} KB"
            else:
                size_str = f"{att_size / (1024 * 1024):.1f} MB"

            # Check size limit (3 MB = 3145728 bytes)
            if att_size > 3 * 1024 * 1024:
                return (
                    f"# Attachment Too Large\n\n"
                    f"**Name:** {att_name}\n"
                    f"**Size:** {size_str}\n"
                    f"**Type:** {att_type}\n\n"
                    f"❌ **Error:** This attachment exceeds the 3 MB limit imposed by Microsoft Graph API.\n\n"
                    f"**Solution:** Please download this attachment manually via Outlook web or desktop app."
                )

            # Download attachment binary
            logger.info(f"Downloading attachment: {att_name} ({size_str})")
            binary_content = await client.download_attachment_binary(
                user_email=user_email,
                message_id=message_id,
                attachment_id=attachment_id,
            )

            logger.info(f"Downloaded {len(binary_content)} bytes")

            # Get file extension
            file_extension = Path(att_name).suffix.lower()

            # Save to temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as temp_file:
                temp_file.write(binary_content)
                temp_file_path = temp_file.name

            try:
                # Extract text based on file type
                extracted_text = ""

                if file_extension == ".pdf":
                    extracted_text = extract_text_from_pdf(temp_file_path)
                elif file_extension == ".docx":
                    extracted_text = extract_text_from_docx(temp_file_path)
                elif file_extension in [".xlsx", ".xls"]:
                    extracted_text = extract_text_from_xlsx(temp_file_path)
                elif file_extension in [".pptx", ".ppt"]:
                    extracted_text = extract_text_from_pptx(temp_file_path)
                elif file_extension in [".html", ".htm"]:
                    extracted_text = extract_text_from_html(temp_file_path)
                elif file_extension in [".txt", ".md", ".csv", ".eml"]:
                    with open(temp_file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        extracted_text = f.read()
                else:
                    return (
                        f"# Unsupported Format\n\n"
                        f"**Name:** {att_name}\n"
                        f"**Size:** {size_str}\n"
                        f"**Type:** {att_type}\n"
                        f"**Format:** {file_extension}\n\n"
                        f"❌ **Error:** This file format is not currently supported.\n\n"
                        f"**Supported formats:**\n"
                        f"- PDF (`.pdf`)\n"
                        f"- Word (`.docx`)\n"
                        f"- Excel (`.xlsx`)\n"
                        f"- PowerPoint (`.pptx`)\n"
                        f"- HTML (`.html`, `.htm`)\n"
                        f"- Text (`.txt`, `.md`, `.csv`, `.eml`)"
                    )

                logger.info(f"Text extraction completed: {len(extracted_text)} characters")

                if not extracted_text or len(extracted_text.strip()) == 0:
                    return (
                        f"# Empty Document Content\n\n"
                        f"**Original File:** {att_name}\n"
                        f"**Size:** {size_str}\n"
                        f"**Type:** {att_type}\n\n"
                        f"⚠️ **Warning:** The document was processed successfully, but no text content was extracted.\n\n"
                        f"This could mean:\n"
                        f"- The document is empty or contains only images/diagrams\n"
                        f"- The document uses unsupported formatting\n"
                        f"- For PDFs: The text may be embedded as images (OCR not available)"
                    )

                # Apply pagination
                truncated = truncate_text_with_metadata(extracted_text, start_word, max_words)

                # Build result
                result_parts = [f"# Extracted Content: {att_name}\n"]
                result_parts.append(f"**Original File:** {att_name}")
                result_parts.append(f"**Size:** {size_str}")
                result_parts.append(f"**Type:** {att_type}")
                result_parts.append(f"**Total Words:** {truncated['total_words']}")

                if truncated["is_truncated"]:
                    result_parts.append(
                        f"**Showing:** Words {truncated['start_word']}-{truncated['end_word']} "
                        f"(of {truncated['total_words']})"
                    )
                    result_parts.append("")
                    result_parts.append(
                        f"*To continue reading, use: `start_word={truncated['end_word']}`*\n"
                    )
                else:
                    result_parts.append(f"**Showing:** Complete document\n")

                result_parts.append("---\n")
                result_parts.append(truncated["text"])

                return "\n".join(result_parts)

            finally:
                # Clean up temp file
                if os.path.exists(temp_file_path):
                    try:
                        os.unlink(temp_file_path)
                        logger.info(f"Cleaned up temp file: {temp_file_path}")
                    except Exception as cleanup_error:
                        logger.warning(f"Failed to clean up temp file: {cleanup_error}")

        except Exception as e:
            error_msg = handle_outlook_error(e)
            logger.error(f"Error in download_outlook_attachment: {error_msg}")
            raise ToolExecutionError("download_outlook_attachment", error_msg)


class CreateEmailTool(CustomTool):
    """Create a new email draft or send immediately."""

    toolkit_name = "outlook"
    toolkit_display_name = "E18 Outlook"

    @property
    def name(self) -> str:
        return "create_email"

    @property
    def description(self) -> str:
        return (
            "Create a new Outlook email draft or send immediately. "
            "\n\n"
            "By default, creates a DRAFT that you can review before sending. "
            "Set send_immediately=true to send without creating a draft. "
            "\n\n"
            "BODY FORMAT:\n"
            "- Accepts markdown or plain text in the 'body' field\n"
            "- Automatically converts markdown to HTML for rich formatting\n"
            "- Use markdown for bold (**text**), italic (*text*), lists, links, etc.\n"
            "\n"
            "Returns the draft message ID (for editing) or confirmation if sent."
        )

    def get_parameters(self) -> List[ToolParameter]:
        return [
            ToolParameter(
                name="to_recipients",
                type="array",
                description="List of recipient email addresses for the 'To' field.",
                required=True,
                items={"type": "string"},
            ),
            ToolParameter(
                name="subject",
                type="string",
                description="Email subject line.",
                required=True,
            ),
            ToolParameter(
                name="body",
                type="string",
                description="Email body content. Supports markdown formatting (converted to HTML automatically).",
                required=True,
            ),
            ToolParameter(
                name="cc_recipients",
                type="array",
                description="List of CC recipient email addresses.",
                required=False,
                items={"type": "string"},
            ),
            ToolParameter(
                name="bcc_recipients",
                type="array",
                description="List of BCC recipient email addresses.",
                required=False,
                items={"type": "string"},
            ),
            ToolParameter(
                name="importance",
                type="string",
                description="Email importance level.",
                required=False,
                enum=["low", "normal", "high"],
                default="normal",
            ),
            ToolParameter(
                name="send_immediately",
                type="boolean",
                description="If true, send immediately. If false (default), create as draft.",
                required=False,
                default=False,
            ),
        ]

    async def _execute_impl(self, user_id: str, **kwargs: Any) -> str:
        """Execute create email tool."""
        try:
            user_email = kwargs.get("_context_user_email")

            if not user_email:
                raise ToolExecutionError(
                    "create_email",
                    "User email is required."
                )

            to_recipients = kwargs.get("to_recipients", [])
            subject = kwargs.get("subject")
            body_markdown = kwargs.get("body")
            cc_recipients = kwargs.get("cc_recipients")
            bcc_recipients = kwargs.get("bcc_recipients")
            importance = kwargs.get("importance", "normal")
            send_immediately = kwargs.get("send_immediately", False)

            if not to_recipients:
                raise ToolExecutionError("create_email", "to_recipients is required")
            if not subject:
                raise ToolExecutionError("create_email", "subject is required")
            if not body_markdown:
                raise ToolExecutionError("create_email", "body is required")

            # Convert markdown to HTML
            body_html = markdown_to_html(body_markdown)

            # Create draft
            client = get_outlook_client()
            draft = await client.create_draft(
                user_email=user_email,
                subject=subject,
                body_content=body_html,
                body_format="HTML",
                to_recipients=to_recipients,
                cc_recipients=cc_recipients,
                bcc_recipients=bcc_recipients,
                importance=importance,
            )

            draft_id = draft.get("id")

            # Send if requested
            if send_immediately:
                await client.send_draft(user_email=user_email, message_id=draft_id)

                return (
                    f"# Email Sent Successfully\n\n"
                    f"**Subject:** {subject}\n"
                    f"**To:** {', '.join(to_recipients)}\n"
                    f"**Status:** Sent and moved to Sent Items folder\n"
                )
            else:
                return (
                    f"# Draft Email Created\n\n"
                    f"**Subject:** {subject}\n"
                    f"**To:** {', '.join(to_recipients)}\n"
                    f"**Draft ID:** `{draft_id}`\n\n"
                    f"*Draft saved to Drafts folder. Use edit_email_draft to modify or send_email_draft to send.*"
                )

        except Exception as e:
            error_msg = handle_outlook_error(e)
            logger.error(f"Error in create_email: {error_msg}")
            raise ToolExecutionError("create_email", error_msg)


class EditEmailDraftTool(CustomTool):
    """Edit an existing email draft."""

    toolkit_name = "outlook"
    toolkit_display_name = "E18 Outlook"

    @property
    def name(self) -> str:
        return "edit_email_draft"

    @property
    def description(self) -> str:
        return (
            "Edit an existing Outlook email draft. "
            "\n\n"
            "Can update subject, body, recipients, and importance. "
            "Only works on draft messages (not sent emails). "
            "\n\n"
            "Use the draft ID from create_email tool."
        )

    def get_parameters(self) -> List[ToolParameter]:
        return [
            ToolParameter(
                name="message_id",
                type="string",
                description="The draft message ID from create_email.",
                required=True,
            ),
            ToolParameter(
                name="subject",
                type="string",
                description="Updated subject line (omit to keep unchanged).",
                required=False,
            ),
            ToolParameter(
                name="body",
                type="string",
                description="Updated body content in markdown (omit to keep unchanged).",
                required=False,
            ),
            ToolParameter(
                name="to_recipients",
                type="array",
                description="Updated 'To' recipients (omit to keep unchanged, or provide array to replace all recipients).",
                required=False,
                items={"type": "string"},
            ),
            ToolParameter(
                name="cc_recipients",
                type="array",
                description="Updated CC recipients (omit to keep unchanged, or provide array to replace all CC recipients).",
                required=False,
                items={"type": "string"},
            ),
            ToolParameter(
                name="bcc_recipients",
                type="array",
                description="Updated BCC recipients (omit to keep unchanged, or provide array to replace all BCC recipients).",
                required=False,
                items={"type": "string"},
            ),
            ToolParameter(
                name="importance",
                type="string",
                description="Updated importance level.",
                required=False,
                enum=["low", "normal", "high"],
            ),
        ]

    async def _execute_impl(self, user_id: str, **kwargs: Any) -> str:
        """Execute edit draft tool."""
        try:
            user_email = kwargs.get("_context_user_email")

            if not user_email:
                raise ToolExecutionError("edit_email_draft", "User email is required.")

            message_id = kwargs.get("message_id")
            if not message_id:
                raise ToolExecutionError("edit_email_draft", "message_id is required")

            # Get update fields
            subject = kwargs.get("subject")
            body_markdown = kwargs.get("body")
            to_recipients = kwargs.get("to_recipients")
            cc_recipients = kwargs.get("cc_recipients")
            bcc_recipients = kwargs.get("bcc_recipients")
            importance = kwargs.get("importance")

            # Convert body if provided
            body_html = None
            if body_markdown:
                body_html = markdown_to_html(body_markdown)

            # Update draft
            client = get_outlook_client()
            updated = await client.update_draft(
                user_email=user_email,
                message_id=message_id,
                subject=subject,
                body_content=body_html,
                body_format="HTML" if body_html else None,
                to_recipients=to_recipients,
                cc_recipients=cc_recipients,
                bcc_recipients=bcc_recipients,
                importance=importance,
            )

            return (
                f"# Draft Updated Successfully\n\n"
                f"**Message ID:** `{message_id}`\n"
                f"**Subject:** {updated.get('subject', 'N/A')}\n\n"
                f"*Draft has been updated. Use send_email_draft to send it.*"
            )

        except Exception as e:
            error_msg = handle_outlook_error(e)
            logger.error(f"Error in edit_email_draft: {error_msg}")
            raise ToolExecutionError("edit_email_draft", error_msg)


class SendEmailDraftTool(CustomTool):
    """Send an existing email draft."""

    toolkit_name = "outlook"
    toolkit_display_name = "E18 Outlook"

    @property
    def name(self) -> str:
        return "send_email_draft"

    @property
    def description(self) -> str:
        return (
            "Send an existing Outlook email draft. "
            "\n\n"
            "The draft will be sent and moved to the Sent Items folder. "
            "This action cannot be undone. "
            "\n\n"
            "Use the draft ID from create_email or edit_email_draft tools."
        )

    def get_parameters(self) -> List[ToolParameter]:
        return [
            ToolParameter(
                name="message_id",
                type="string",
                description="The draft message ID to send.",
                required=True,
            ),
        ]

    async def _execute_impl(self, user_id: str, **kwargs: Any) -> str:
        """Execute send draft tool."""
        try:
            user_email = kwargs.get("_context_user_email")

            if not user_email:
                raise ToolExecutionError("send_email_draft", "User email is required.")

            message_id = kwargs.get("message_id")
            if not message_id:
                raise ToolExecutionError("send_email_draft", "message_id is required")

            # Send draft
            client = get_outlook_client()
            await client.send_draft(user_email=user_email, message_id=message_id)

            return (
                f"# Email Sent Successfully\n\n"
                f"**Message ID:** `{message_id}`\n"
                f"**Status:** Sent and moved to Sent Items folder\n"
            )

        except Exception as e:
            error_msg = handle_outlook_error(e)
            logger.error(f"Error in send_email_draft: {error_msg}")
            raise ToolExecutionError("send_email_draft", error_msg)
